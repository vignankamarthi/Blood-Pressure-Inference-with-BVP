#!/usr/bin/env python3
"""
DS4400 Final Project Presentation Generator
Blood Pressure Inference from Physiological Signals Using
Feature Engineering and Deep Learning

Generates a 5-slide .pptx with dark theme, embedded figures, and custom charts.
Run: python3 generate_presentation.py
"""

import os
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# Configuration
# ============================================================

BASE = os.path.dirname(os.path.abspath(__file__))
FIG_DIR = os.path.join(BASE, "figures")
os.makedirs(FIG_DIR, exist_ok=True)

# Color palette
BG_HEX = '#0F1624'
BG_RGB = RGBColor(0x0F, 0x16, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SILVER = RGBColor(0xB0, 0xB0, 0xB0)
LGRAY = RGBColor(0xD0, 0xD0, 0xD0)
DIM = RGBColor(0x55, 0x55, 0x55)
NEU_RED = RGBColor(0xC8, 0x10, 0x2E)
GOLD = RGBColor(0xFF, 0xD7, 0x00)
TEAL = RGBColor(0x4E, 0xCD, 0xC4)

TOTAL_SLIDES = 5

# Matplotlib dark theme
plt.rcParams.update({
    'figure.facecolor': BG_HEX,
    'axes.facecolor': BG_HEX,
    'text.color': 'white',
    'axes.labelcolor': '#999999',
    'xtick.color': '#999999',
    'ytick.color': '#999999',
    'axes.edgecolor': '#444444',
    'font.family': 'sans-serif',
})


# ============================================================
# Helpers
# ============================================================

def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_RGB
    return slide


def _set_run_font(run, size, color, bold=False, italic=False, name='Calibri'):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def add_text(slide, text, left, top, width, height,
             size=20, color=SILVER, bold=False, italic=False,
             align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if p.runs:
            _set_run_font(p.runs[0], size, color, bold, italic)
    return tf


def add_title(slide, text, left=Inches(0.8), top=Inches(0.3)):
    tf = add_text(slide, text, left, top, Inches(11.5), Inches(0.9),
                  size=36, color=WHITE, bold=True)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.95), Inches(1.5), Pt(4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = NEU_RED
    accent.line.color.rgb = NEU_RED
    accent.line.width = Pt(0)
    return tf


def add_bullets(slide, items, left, top, width, height, size=22):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, color) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {text}"
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(16)
        if p.runs:
            _set_run_font(p.runs[0], size, color or SILVER)
    return tf


def add_image(slide, path, left, top, width=None, height=None):
    if not os.path.exists(path):
        print(f"  WARNING: missing {path}")
        return None
    kw = {'image_file': path, 'left': left, 'top': top}
    if width:
        kw['width'] = width
    if height:
        kw['height'] = height
    return slide.shapes.add_picture(**kw)


def add_footer(slide, num):
    add_text(slide, "DS 4400  |  Spring 2026  |  Northeastern University",
             Inches(0.5), Inches(7.05), Inches(4), Inches(0.35),
             size=9, color=DIM)
    add_text(slide, f"{num} / {TOTAL_SLIDES}",
             Inches(11.8), Inches(7.05), Inches(1.2), Inches(0.35),
             size=9, color=DIM, align=PP_ALIGN.RIGHT)


# ============================================================
# Figure Generation
# ============================================================

def create_pipeline_diagram():
    """Two-track pipeline flow diagram: classical ML vs deep learning."""
    fig, ax = plt.subplots(figsize=(6.5, 4.5))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 8)
    ax.axis('off')

    # Box drawing helper
    def draw_box(x, y, w, h, text, facecolor, edgecolor, fontsize=11,
                 textcolor='white', fontweight='normal'):
        rect = plt.Rectangle((x, y), w, h, facecolor=facecolor,
                              edgecolor=edgecolor, linewidth=1.5,
                              joinstyle='round')
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, text, ha='center', va='center',
                fontsize=fontsize, color=textcolor, fontweight=fontweight,
                wrap=True)

    # Arrow helper
    def draw_arrow(x1, y1, x2, y2, color='#999999'):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle='->', color=color, lw=1.8))

    # Title labels
    ax.text(2.5, 7.5, 'Classical ML Track', fontsize=13, color='#4ECDC4',
            ha='center', fontweight='bold')
    ax.text(7.5, 7.5, 'Deep Learning Track', fontsize=13, color='#FFD700',
            ha='center', fontweight='bold')

    # Vertical separator
    ax.plot([5, 5], [0.5, 7.2], color='#444444', lw=1, ls='--')

    # --- Classical ML Track (left) ---
    # Raw signals
    draw_box(1, 6.2, 3, 0.8, 'Raw PPG + ECG + ABP', '#1a2744', '#4ECDC4')
    draw_arrow(2.5, 6.2, 2.5, 5.6, '#4ECDC4')

    # Rust feature extraction
    draw_box(0.8, 4.8, 3.4, 0.8, 'Rust Feature Extraction\n(40 feat/signal)',
             '#1a2744', '#4ECDC4', fontsize=10)
    draw_arrow(2.5, 4.8, 2.5, 4.2, '#4ECDC4')

    # Optuna tuning
    draw_box(1, 3.4, 3, 0.8, 'Optuna Tuning\n(30 trials/model)',
             '#1a2744', '#4ECDC4', fontsize=10)
    draw_arrow(2.5, 3.4, 2.5, 2.8, '#4ECDC4')

    # 5 ML models
    draw_box(0.5, 2.0, 4, 0.8, 'Ridge / DT / RF / XGB / LGBM',
             '#1a2744', '#4ECDC4', fontsize=10)

    # --- Deep Learning Track (right) ---
    # Raw PPG only
    draw_box(6, 6.2, 3, 0.8, 'Raw PPG (1ch, 1250)', '#1a2744', '#FFD700')
    draw_arrow(7.5, 6.2, 7.5, 5.6, '#FFD700')

    # ResNet blocks
    draw_box(6, 4.8, 3, 0.8, 'ResNet Residual Blocks',
             '#1a2744', '#FFD700', fontsize=10)
    draw_arrow(7.5, 4.8, 7.5, 4.2, '#FFD700')

    # BiGRU
    draw_box(6, 3.4, 3, 0.8, 'Bidirectional GRU',
             '#1a2744', '#FFD700', fontsize=10)
    draw_arrow(7.5, 3.4, 7.5, 2.8, '#FFD700')

    # DL output
    draw_box(6, 2.0, 3, 0.8, 'ResNet-BiGRU / ResNet-1D',
             '#1a2744', '#FFD700', fontsize=10)

    # --- Both converge ---
    draw_arrow(2.5, 2.0, 4.5, 1.2, '#999999')
    draw_arrow(7.5, 2.0, 5.5, 1.2, '#999999')

    draw_box(3.5, 0.4, 3, 0.8, 'SBP + DBP\nEvaluation',
             '#C8102E', '#C8102E', fontsize=12, fontweight='bold')

    plt.tight_layout()
    path = os.path.join(FIG_DIR, 'pipeline.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=BG_HEX)
    plt.close()
    print(f"  Created: {os.path.basename(path)}")
    return path


def create_results_comparison():
    """Grouped bar chart: LGBM vs ResNet-BiGRU on SBP and DBP MAE."""
    fig, ax = plt.subplots(figsize=(6.5, 4.2))

    metrics = ['SBP MAE', 'DBP MAE']
    lgbm_vals = [14.46, 8.54]
    resnet_vals = [13.61, 7.97]

    x = np.arange(len(metrics))
    bar_w = 0.3

    bars1 = ax.bar(x - bar_w / 2, lgbm_vals, bar_w, label='LightGBM (Classical)',
                   color='#555555', edgecolor='none')
    bars2 = ax.bar(x + bar_w / 2, resnet_vals, bar_w, label='ResNet-BiGRU (DL)',
                   color='#C8102E', edgecolor='none')

    # Value labels on bars
    for bar, val in zip(bars1, lgbm_vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                f'{val}', ha='center', color='#B0B0B0', fontsize=14,
                fontweight='bold')

    for bar, val in zip(bars2, resnet_vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.3,
                f'{val}', ha='center', color='#FFD700', fontsize=14,
                fontweight='bold')

    # Improvement arrows
    for i in range(len(metrics)):
        diff = lgbm_vals[i] - resnet_vals[i]
        mid_x = x[i] + bar_w / 2
        ax.annotate(f'-{diff:.2f}',
                    xy=(mid_x, resnet_vals[i] + 0.15),
                    xytext=(mid_x + 0.45, resnet_vals[i] + 2.5),
                    fontsize=10, color='#4ECDC4', fontweight='bold',
                    arrowprops=dict(arrowstyle='->', color='#4ECDC4', lw=1.2))

    ax.set_xticks(x)
    ax.set_xticklabels(metrics, fontsize=14, color='white')
    ax.set_ylabel('MAE (mmHg)', fontsize=13)
    ax.set_ylim(0, 18)
    ax.legend(fontsize=11, loc='upper right', bbox_to_anchor=(1.0, 1.15),
              framealpha=0.3, facecolor=BG_HEX, edgecolor='#444444',
              labelcolor='white', ncol=2)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)

    plt.tight_layout()
    path = os.path.join(FIG_DIR, 'results_comparison.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=BG_HEX)
    plt.close()
    print(f"  Created: {os.path.basename(path)}")
    return path


# ============================================================
# Slide Builders
# ============================================================

def slide_1_title(prs):
    """Title slide (~30s) [Ariv speaks]"""
    s = new_slide(prs)

    add_text(s, "CUFFLESS BLOOD PRESSURE ESTIMATION",
             Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.2),
             size=44, color=WHITE, bold=True)

    add_text(s, "From Physiological Signals Using\nFeature Engineering and Deep Learning",
             Inches(0.8), Inches(2.6), Inches(11.5), Inches(1),
             size=28, color=TEAL, italic=True)

    # Accent line
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.8), Inches(3.7), Inches(3), Pt(4))
    acc.fill.solid()
    acc.fill.fore_color.rgb = NEU_RED
    acc.line.color.rgb = NEU_RED

    add_text(s, "Vignan Kamarthi  &  Ariv Ahuja",
             Inches(0.8), Inches(4.1), Inches(6), Inches(0.6),
             size=26, color=WHITE, bold=True)

    add_text(s, "DS 4400: Machine Learning  |  Prof. Silvio Amir  |  Spring 2026",
             Inches(0.8), Inches(4.9), Inches(11), Inches(0.5),
             size=16, color=DIM)

    add_text(s, "github.com/vignankamarthi/Blood-Pressure-Inference-with-BVP",
             Inches(0.8), Inches(6.3), Inches(11), Inches(0.4),
             size=14, color=DIM)


def slide_2_problem_dataset(prs):
    """Problem + Dataset slide (~60s) [Ariv speaks]"""
    s = new_slide(prs)
    add_title(s, "The Clinical Challenge")

    # Left side: clinical motivation
    add_text(s, "140/90 mmHg",
             Inches(1.0), Inches(1.8), Inches(5), Inches(0.7),
             size=36, color=GOLD, bold=True)

    items_left = [
        ("Hypertension threshold: leading cause of cardiovascular death", SILVER),
        ("Cuff-based measurement: infrequent, clinic-only, reactive", LGRAY),
        ("PPG from wearables: continuous, passive, preventive", TEAL),
    ]
    add_bullets(s, items_left, Inches(1.0), Inches(2.8), Inches(5.5), Inches(3), size=20)

    # Right side: dataset stats
    add_text(s, "PulseDB v2.0",
             Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.6),
             size=24, color=WHITE, bold=True)

    add_text(s, "MIMIC + VitalDB",
             Inches(7.2), Inches(2.4), Inches(5.5), Inches(0.4),
             size=16, color=DIM)

    stats = [
        ("5.2M segments, 5,361 subjects", GOLD),
        ("3 signals: PPG, ECG, ABP", SILVER),
        ("125 Hz, 10-second windows", SILVER),
        ("CalFree protocol: train/test by subject", LGRAY),
    ]
    add_bullets(s, stats, Inches(7.2), Inches(3.1), Inches(5.5), Inches(3), size=18)

    # Bottom teal callout
    add_text(s, "Two-track experiment: handcrafted features vs. learned representations",
             Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
             size=18, color=TEAL, italic=True, align=PP_ALIGN.CENTER)

    add_footer(s, 2)


def slide_3_methods(prs):
    """Methods slide (~90s) [Vignan speaks]"""
    s = new_slide(prs)
    add_title(s, "Two-Track Pipeline")

    # Pipeline diagram (left)
    pipeline_fig = os.path.join(FIG_DIR, 'pipeline.png')
    add_image(s, pipeline_fig, Inches(0.2), Inches(1.5), width=Inches(6.5))

    # Right side: method details
    add_text(s, "Feature Engineering",
             Inches(7.2), Inches(1.6), Inches(5.5), Inches(0.5),
             size=20, color=TEAL, bold=True)

    items_fe = [
        ("Catch22 (22) + Entropy (10) + Stats (8)", SILVER),
        ("Rust extraction: 40 features per signal", SILVER),
        ("Optuna tuning: 30 trials per model", LGRAY),
    ]
    add_bullets(s, items_fe, Inches(7.2), Inches(2.2), Inches(5.5), Inches(2), size=17)

    add_text(s, "Deep Learning",
             Inches(7.2), Inches(4.0), Inches(5.5), Inches(0.5),
             size=20, color=GOLD, bold=True)

    items_dl = [
        ("ResNet-BiGRU: residual blocks + bidirectional GRU", SILVER),
        ("ResNet-1D baseline for comparison", SILVER),
        ("GradientSHAP temporal attribution", LGRAY),
    ]
    add_bullets(s, items_dl, Inches(7.2), Inches(4.6), Inches(5.5), Inches(2), size=17)

    add_footer(s, 3)


def slide_4_results(prs):
    """Results slide (~90s) [Vignan speaks]"""
    s = new_slide(prs)
    add_title(s, "Key Results (CalFree Test Set)")

    # Results comparison chart (left)
    chart = os.path.join(FIG_DIR, 'results_comparison.png')
    add_image(s, chart, Inches(0.2), Inches(1.5), width=Inches(6.5))

    # Right side: big numbers
    add_text(s, "13.61",
             Inches(7.2), Inches(1.5), Inches(5.5), Inches(1.2),
             size=60, color=GOLD, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, "SBP MAE (mmHg)",
             Inches(7.2), Inches(2.7), Inches(5.5), Inches(0.5),
             size=18, color=SILVER, align=PP_ALIGN.CENTER)

    add_text(s, "BHS Grade C",
             Inches(7.2), Inches(3.5), Inches(5.5), Inches(0.6),
             size=28, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    add_text(s, "for DBP (best model)",
             Inches(7.2), Inches(4.1), Inches(5.5), Inches(0.4),
             size=14, color=DIM, align=PP_ALIGN.CENTER)

    # Bottom bullets
    bottom_items = [
        ("ECG adds only 0.03 mmHg improvement over PPG alone", SILVER),
        ("All models fail AAMI standard (SD > 8 mmHg)", SILVER),
        ("Consistent with PulseDB benchmark results", SILVER),
    ]
    add_bullets(s, bottom_items, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8),
                size=17)

    add_footer(s, 4)


def slide_5_conclusions(prs):
    """Conclusions slide (~30s) [Ariv speaks]"""
    s = new_slide(prs)
    add_title(s, "What We Learned")

    # Three key takeaways as large items
    takeaways = [
        ("PPG alone is sufficient for wearable BP monitoring", TEAL),
        ("Deep learning outperforms feature engineering by 0.85 mmHg", GOLD),
        ("AAMI compliance remains an open challenge for the field", SILVER),
    ]

    y_pos = Inches(2.0)
    for text, color in takeaways:
        add_text(s, f"\u2022   {text}",
                 Inches(1.5), y_pos, Inches(10), Inches(0.8),
                 size=26, color=color, bold=False)
        y_pos += Inches(1.2)

    # Thank you
    add_text(s, "Thank you",
             Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7),
             size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Contact info
    add_text(s, "kamarthi.v@northeastern.edu  |  ahuja.ar@northeastern.edu",
             Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             size=14, color=DIM, align=PP_ALIGN.CENTER)

    add_footer(s, 5)


# ============================================================
# Main
# ============================================================

def main():
    print("=== DS4400 Final Project Presentation Generator ===\n")

    print("Generating figures...")
    create_pipeline_diagram()
    create_results_comparison()

    print("\nBuilding presentation (5 slides)...")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_1_title(prs)
    print("  Slide 1: Title")
    slide_2_problem_dataset(prs)
    print("  Slide 2: Problem + Dataset")
    slide_3_methods(prs)
    print("  Slide 3: Methods")
    slide_4_results(prs)
    print("  Slide 4: Results")
    slide_5_conclusions(prs)
    print("  Slide 5: Conclusions")

    out_path = os.path.join(BASE, "DS4400_Final_Presentation.pptx")
    prs.save(out_path)
    print(f"\nSaved: {out_path}")
    return out_path


if __name__ == "__main__":
    main()
